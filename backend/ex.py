import os
import re
import tempfile
from io import BytesIO
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import pyttsx3
import google.generativeai as genai
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
from werkzeug.utils import secure_filename
import traceback

import torch
from diffusers import StableDiffusionPipeline
import textwrap
import math
# from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
from PIL import Image, ImageDraw, ImageFont
from gradio_client import Client, handle_file

from flask import Flask, request, jsonify
from flask_cors import CORS
from flask_bcrypt import Bcrypt
import jwt
import datetime

## For Reddis Caching
import redis
import hashlib
import pickle

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}}, supports_credentials=True)

redis_client = redis.Redis(host='localhost', port=6379, db=0, decode_responses=False)

def redis_key_for_text(text):
    return "embed:" + hashlib.sha256(text.encode('utf-8')).hexdigest()

bcrypt = Bcrypt(app)
app.config['SECRET_KEY'] = 'abc@123'  # Change this!

# In-memory user store (replace with DB in production)
from pymongo import MongoClient
import os
MONGO_URI =os.getenv("MONGO_URI","mongodb://localhost:27017")

mongo_client = MongoClient(MONGO_URI)
db = mongo_client["researchhive"]
users_collection = db["users"]
users = {}

def generate_token(username):
    payload = {
        'username': username,
        'exp': datetime.datetime.utcnow() + datetime.timedelta(hours=12)
    }
    return jwt.encode(payload, app.config['SECRET_KEY'], algorithm='HS256')

def verify_token(token):
    try:
        payload = jwt.decode(token, app.config['SECRET_KEY'], algorithms=['HS256'])
        return payload['username']
    except Exception:
        return None

@app.route('/register', methods=['POST'])
def register():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    if not username or not password:
        return jsonify({'error': 'Username and password required'}), 400
    if username in users:
        return jsonify({'error': 'User already exists'}), 400
    hashed = bcrypt.generate_password_hash(password).decode('utf-8')
    users_collection.insert_one({"username": username, "password": hashed})
    return jsonify({'message': 'User registered successfully'}), 200

@app.route('/login', methods=['POST'])
def login():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    if not username or not password:
        return jsonify({'error': 'Username and password required'}), 400
    user = users_collection.find_one({"username": username})
    if not user or not bcrypt.check_password_hash(user["password"], password):
        return jsonify({'error': 'Invalid credentials'}), 401
    token = generate_token(username)
    return jsonify({'token': token}), 200

# Example protected route
@app.route('/protected', methods=['GET'])
def protected():
    token = request.headers.get('Authorization', '').replace('Bearer ', '')
    username = verify_token(token)
    if not username:
        return jsonify({'error': 'Unauthorized'}), 401
    return jsonify({'message': f'Hello, {username}!'}), 200

genai.configure(api_key="AIzaSyB3aKZrfQZmfISU-nay6reeFp_At1BgE50")  # Replace with your Gemini API key
model = genai.GenerativeModel("gemini-2.5-flash")

level_prompts = {
    "Beginner": "Summarize this research paper section for a high school student in 4-6 concise bullet points:",
    "Student": "Create a structured summary of this section for undergraduate students in 4-6 points:",
    "Expert": "Generate a detailed summary of this section for researchers in 5-7 well-formed bullet points:"
}

creativity_levels = {
    "Formal": "Keep the conversation strictly professional and formal.",
    "Balanced": "Maintain a balance between professional and conversational tone.",
    "Creative": "Make the conversation more creative and engaging with some informal elements."
}

podcast_lengths = {
    "Short (2-3 mins)": "Generate a short podcast with 2-3 questions and concise answers.",
    "Medium (5-7 mins)": "Generate a medium-length podcast with 4-5 questions and detailed answers.",
    "Long (10+ mins)": "Generate a long podcast with 6-8 questions and in-depth discussion."
}

template_options = {
    "Template 1": "templates/theme_template_1.pptx",
    "Template 2": "templates/theme_template_2.pptx",
    "Template 3": "templates/theme_template_3.pptx"
}

from sentence_transformers import SentenceTransformer
import faiss
import numpy as np

def chunk_text(text, chunk_size=300, overlap=50):
    """Split text into overlapping chunks."""
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i+chunk_size])
        if chunk:
            chunks.append(chunk)
    return chunks

embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

def build_faiss_index_with_cache(chunks):
    """Create a FAISS index from text chunks, caching embeddings in Redis."""
    embeddings = []
    for chunk in chunks:
        key = redis_key_for_text(chunk)
        cached = redis_client.get(key)
        if cached:
            emb = pickle.loads(cached)
        else:
            emb = embedding_model.encode([chunk])[0]
            redis_client.set(key, pickle.dumps(emb))
        embeddings.append(emb)
    embeddings = np.stack(embeddings)
    index = faiss.IndexFlatL2(embeddings.shape[1])
    index.add(embeddings.astype('float32'))
    return index, embeddings, chunks

def retrieve_relevant_chunks_with_cache(query, index, embeddings, chunks, top_k=3):
    """Retrieve top-k relevant chunks for a query, caching query embedding in Redis."""
    key = redis_key_for_text(query)
    cached = redis_client.get(key)
    if cached:
        query_emb = pickle.loads(cached)
    else:
        query_emb = embedding_model.encode([query])[0]
        redis_client.set(key, pickle.dumps(query_emb))
    query_emb = np.array(query_emb).reshape(1, -1).astype('float32')
    D, I = index.search(query_emb, top_k)
    return [chunks[i] for i in I[0]]

def rag_generate_answer(query, document_text):
    
    chunks = chunk_text(document_text)
    
    index, embeddings, chunk_list = build_faiss_index_with_cache(chunks)
    
    retrieved_chunks = retrieve_relevant_chunks_with_cache(query, index, embeddings, chunk_list, top_k=3)
    
    context = "\n".join(retrieved_chunks)
    prompt = f"""Use the following context to answer the question:\n\nContext:\n{context}\n\nQuestion: {query}\nAnswer:"""
    
    response = model.generate_content(prompt)
    return response.text

import time 

@app.route('/rag-answer', methods=['POST'])
def rag_answer():
    try:
        data = request.json
        document_text = data.get('document_text')
        query = data.get('query')
        if not document_text or not query:
            return jsonify({'error': 'document_text and query are required'}), 400
        
        start_total = time.time()
        # --- FAISS timing ---
        start_faiss = time.time()
        chunks = chunk_text(document_text)
        index, embeddings, chunk_list = build_faiss_index_with_cache(chunks)
        retrieved_chunks = retrieve_relevant_chunks_with_cache(query, index, embeddings, chunk_list, top_k=3)
        end_faiss = time.time()
        print(f"FAISS indexing and search took {end_faiss - start_faiss:.2f} seconds")
        # --- RAG pipeline ---
        context = "\n".join(retrieved_chunks)
        prompt = f"""Use the following context to answer the question:\n\nContext:\n{context}\n\nQuestion: {query}\nAnswer:"""
        answer_start = time.time()
        response = model.generate_content(prompt)
        answer = response.text
        answer_end = time.time()
        print(f"RAG pipeline answer: {answer}")
        print(f"RAG answer generation took {answer_end - answer_start:.2f} seconds")
        print(f"Total RAG QA pipeline took {time.time() - start_total:.2f} seconds")
        return jsonify({'answer': answer}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def get_avatar_html(active_speaker):
    """Generate HTML for avatars with active speaker highlighting."""
    return f"""
    <div class="avatar-container">
        <div class="avatar-card {'active-speaker' if active_speaker == 'Alex' else ''}">
            <div class="speaking-indicator"></div>
            <img src="https://img.icons8.com/color/144/000000/circled-user-female-skin-type-5.png" 
                 class="avatar-img">
            <div class="avatar-name">Alex</div>
        </div>
        <div class="avatar-card {'active-speaker' if active_speaker == 'Dr. Smith' else ''}">
            <div class="speaking-indicator"></div>
            <img src="https://img.icons8.com/color/144/000000/circled-user-male-skin-type-7.png" 
                 class="avatar-img">
            <div class="avatar-name">Dr. Smith</div>
        </div>
    </div>
    """

def generate_content_from_heading(heading):
    """Generate research paper-like content from a heading using Gemini."""
    prompt = f"""You are a research paper writer. Based on the following heading, generate a detailed research paper-like content with the following structure:
1. Introduction
2. Methodology
3. Results
4. Discussion
5. Conclusion

Heading: {heading}

Ensure the content is well-structured, informative, and suitable for academic purposes.
"""
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        raise Exception(f"Error generating content: {e}")

def extract_and_summarize_sections(text, summary_level):
    """Extract sections and generate summaries using Gemini."""
    prompt = f"""Analyze the following research paper and:
1. Identify all major sections.
2. For each section, generate a summary using the following guidelines:
   - {level_prompts[summary_level]}
3. Format the response as:
   ## Section Name
   - Bullet point 1
   - Bullet point 2
   - Bullet point 3

Paper content:
{text}
"""
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        raise Exception(f"Error processing document: {e}")

def generate_podcast_script(summary_text, creativity_level, podcast_length):
    """Generate a conversational podcast script using Gemini."""
    prompt = f"""Create a conversational podcast script between host Alex and researcher Dr. Smith discussing the research paper. Follow these rules:
1. Alex should ask curious, layperson-friendly questions.
2. Dr. Smith should provide expert answers based on the paper.
3. Always prefix lines with either "Alex:" or "Dr. Smith:".
4. Keep responses conversational but informative.
5. Cover key findings, methodology, and implications.
6. {creativity_levels[creativity_level]}
7. {podcast_lengths[podcast_length]}

Paper summary:
{summary_text}
"""
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        raise Exception(f"Error generating podcast script: {e}")

def create_ppt_from_summary(summary_text, template_path):
    """Create PowerPoint from section-wise summaries using the selected template."""
    prs = Presentation(template_path)

    # Generate a title for the presentation using Gemini
    title_prompt = f"""Analyze the following text and generate a concise, professional title for a PowerPoint presentation (maximum 10-12 words):
    {summary_text[:5000]}  # Use the first 5000 characters for title generation
    """
    try:
        title_response = model.generate_content(title_prompt)
        title = title_response.text.strip()
    except Exception as e:
        title = "Research Summary"  

    max_title_length = 80  
    if len(title) > max_title_length:
        words = title.split()
        line1 = ""
        line2 = ""
        for word in words:
            if len(line1) + len(word) + 1 <= max_title_length:
                line1 += word + " "
            else:
                line2 += word + " "
        title = f"{line1.strip()}\n{line2.strip()}"

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated by AI Driven Multimedia Generator"

    slide_layout = prs.slide_layouts[1]

    sections = {}
    current_section = "Introduction"
    for line in summary_text.split('\n'):
        if line.startswith("## "):
            current_section = line[3:].strip()
            sections[current_section] = []
        elif line.startswith("- "):
            sections[current_section].append(line[2:])

    for section, bullets in sections.items():
        slides_per_section = min((len(bullets) // 6) + 1, 5)
        chunk_size = max(len(bullets) // slides_per_section, 1)
        
        for i in range(0, len(bullets), chunk_size):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = section if i == 0 else f"{section} (Cont.)"
            
            content_box = slide.shapes.placeholders[1]
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for bullet in bullets[i:i+chunk_size]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = 0
    
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

import tempfile
import os
import pyttsx3
import wave
import struct
import re
import time

def generate_podcast_audio(podcast_script, rate=150):
    """Generate TTS audio with distinct voices for host and researcher using separate engines."""
    try:
        temp_dir = tempfile.mkdtemp()
        output_file = os.path.join(temp_dir, "podcast_audio.wav")
        
        segments = parse_podcast_script(podcast_script)
        
        if not segments:
            raise Exception("No valid speaker segments found in script")
        
        print(f"Found {len(segments)} speaker segments")
        
        # Generate audio for each segment
        audio_files = []
        
        for i, segment in enumerate(segments):
            speaker = segment['speaker']
            text = segment['text']
            
            if not text or len(text.strip()) < 3:
                continue
            
            # Generate audio for this segment
            segment_file = generate_speaker_audio(
                text, 
                speaker, 
                temp_dir, 
                f"segment_{i}", 
                rate
            )
            
            if segment_file and os.path.exists(segment_file):
                audio_files.append(segment_file)
                print(f"Generated audio for {speaker}: {text[:50]}...")
        
        if not audio_files:
            raise Exception("No audio segments were successfully generated")
        
        # Combine all audio files
        combine_audio_files(audio_files, output_file)
        
        # Clean up temporary segment files
        for file in audio_files:
            try:
                os.remove(file)
            except:
                pass
        
        # Verify output
        if not os.path.exists(output_file) or os.path.getsize(output_file) < 1000:
            raise Exception("Final audio file was not created properly")
        
        print(f"Successfully generated podcast audio: {os.path.getsize(output_file)} bytes")
        return output_file
        
    except Exception as e:
        print(f"Error in generate_podcast_audio: {e}")
        raise Exception(f"Error generating audio: {e}")

## into small parts, to extract the useful info by parsing 
def parse_podcast_script(script):
    """Parse podcast script into speaker segments."""
    segments = []
    lines = [line.strip() for line in script.split('\n') if line.strip()]
    
    for line in lines:
        # Skip empty lines or very short lines
        if len(line) < 5:
            continue
        
        speaker = "host"  # first
        text = line
        
        # Check for different speaker patterns
        if ':' in line:
            before_colon = line.split(':', 1)[0].strip()
            after_colon = line.split(':', 1)[1].strip()
            
            # Identify speaker based on keywords
            before_lower = before_colon.lower()
            
            if any(keyword in before_lower for keyword in ['alex', 'host', 'interviewer']):
                speaker = "host"
                text = after_colon
            elif any(keyword in before_lower for keyword in ['dr.', 'doctor', 'smith', 'researcher', 'expert', 'scientist']):
                speaker = "expert"
                text = after_colon
            elif any(keyword in before_lower for keyword in ['narrator', 'introduction', 'conclusion']):
                speaker = "narrator"
                text = after_colon
            else:
                speaker = "narrator"
                text = line
        
        # Clean up text
        text = text.strip()
        if text and len(text) > 2:
            segments.append({
                'speaker': speaker,
                'text': text
            })
    
    return segments


def generate_speaker_audio(text, speaker, temp_dir, filename, rate=150):
    """Generate audio for a specific speaker with appropriate voice settings."""
    try:
        output_file = os.path.join(temp_dir, f"{filename}.wav")
        
        # Initialize TTS engine
        engine = pyttsx3.init()
        
        # Get available voices
        voices = engine.getProperty('voices')
        
        if not voices:
            raise Exception("No TTS voices available")
        
        # Configure voice based on speaker
        if speaker == "host":
            # Use female voice with moderate rate
            if len(voices) > 1:
                engine.setProperty('voice', voices[1].id)  # Usually female
            else:
                engine.setProperty('voice', voices[0].id)
            engine.setProperty('rate', rate)
            engine.setProperty('volume', 0.9)
            
        elif speaker == "expert":
            # Use male voice with slightly slower rate
            engine.setProperty('voice', voices[0].id)  # Usually male
            engine.setProperty('rate', rate - 20)  # Slightly slower for authority
            engine.setProperty('volume', 0.95)
            
        else:  # narrator
            # Use neutral voice
            engine.setProperty('voice', voices[0].id)
            engine.setProperty('rate', rate + 10)  # Slightly faster for narration
            engine.setProperty('volume', 0.85)
        
        # Generate the audio
        engine.save_to_file(text, output_file)
        engine.runAndWait()
        
        # Give a small delay to ensure file is written
        time.sleep(0.1)
        
        # Verify the file was created
        if os.path.exists(output_file) and os.path.getsize(output_file) > 100:
            return output_file
        else:
            print(f"Failed to generate audio for: {text[:30]}...")
            return None
            
    except Exception as e:
        print(f"Error generating audio for {speaker}: {e}")
        return None


def combine_audio_files(audio_files, output_file):
    """Combine multiple audio files with natural pauses between speakers."""
    try:
        with wave.open(output_file, 'wb') as output_wav:
            params_set = False
            
            for i, audio_file in enumerate(audio_files):
                try:
                    with wave.open(audio_file, 'rb') as input_wav:
                        params = input_wav.getparams()
                        
                        if not params_set:
                            output_wav.setparams(params)
                            params_set = True
                        
                        # Write the audio data
                        frames = input_wav.readframes(input_wav.getnframes())
                        output_wav.writeframes(frames)
                        
                        # Add pause between segments (except for the last one)
                        if i < len(audio_files) - 1:
                            add_silence_to_wav(output_wav, 0.8, params)  # 0.8 second pause
                            
                except Exception as e:
                    print(f"Error processing audio file {audio_file}: {e}")
                    continue
        
        print(f"Successfully combined {len(audio_files)} audio segments")
        
    except Exception as e:
        raise Exception(f"Error combining audio files: {e}")


def add_silence_to_wav(wav_file, duration_seconds, params):
    """Add silence to an open WAV file."""
    try:
        sample_rate = params.framerate
        channels = params.nchannels
        sample_width = params.sampwidth
        
        # Calculate number of frames for silence
        silence_frames = int(duration_seconds * sample_rate)
        
        # Create silence data (zeros)
        silence_data = b'\x00' * (silence_frames * channels * sample_width)
        
        # Write silence
        wav_file.writeframes(silence_data)
        
    except Exception as e:
        print(f"Error adding silence: {e}")


def generate_simple_podcast_audio(podcast_script, rate=150):
    """Simplified fallback version if the main function fails."""
    try:
        temp_dir = tempfile.mkdtemp()
        output_file = os.path.join(temp_dir, "podcast_audio.wav")
        
        # Clean the script
        clean_script = clean_script_for_tts(podcast_script)
        
        if not clean_script:
            raise Exception("No valid content found in script")
        
        # Generate audio with single voice
        engine = pyttsx3.init()
        engine.setProperty('rate', rate)
        
        voices = engine.getProperty('voices')
        if voices:
            engine.setProperty('voice', voices[0].id)
        
        engine.save_to_file(clean_script, output_file)
        engine.runAndWait()
        
        # Verify output
        if not os.path.exists(output_file) or os.path.getsize(output_file) < 500:
            raise Exception("Simple audio generation failed")
        
        return output_file
        
    except Exception as e:
        raise Exception(f"Error in simple audio generation: {e}")


def clean_script_for_tts(script):
    """Clean the script for better TTS output."""
    lines = script.split('\n')
    cleaned_lines = []
    
    for line in lines:
        line = line.strip()
        if not line or len(line) < 3:
            continue
        
        # Remove speaker names and formatting
        if ':' in line:
            text = line.split(':', 1)[1].strip()
        else:
            text = line
        
        # Clean up text
        text = re.sub(r'\s+', ' ', text)  # Multiple spaces to single space
        text = re.sub(r'[^\w\s\.,!?;:-]', '', text)  # Remove special characters
        
        if text and len(text) > 2:
            cleaned_lines.append(text)
    
    # Join with natural pauses
    return ' ... '.join(cleaned_lines)


# Alternative approach using separate audio streams
def generate_dual_voice_podcast(podcast_script, rate=150):
    """Generate podcast with clearly distinct voices using separate processing."""
    try:
        temp_dir = tempfile.mkdtemp()
        output_file = os.path.join(temp_dir, "podcast_audio.wav")
        
        # Separate script by speakers
        host_lines, expert_lines, timeline = separate_speakers(podcast_script)
        
        # Generate separate audio files
        host_audio = generate_voice_audio(host_lines, "host", temp_dir, rate)
        expert_audio = generate_voice_audio(expert_lines, "expert", temp_dir, rate)
        
        # Interleave based on timeline
        interleave_audio_by_timeline(timeline, host_audio, expert_audio, output_file)
        
        # Cleanup
        cleanup_temp_files([host_audio, expert_audio])
        
        return output_file
        
    except Exception as e:
        raise Exception(f"Error in dual voice generation: {e}")


def separate_speakers(script):
    """Separate script into host and expert lines with timeline."""
    host_lines = []
    expert_lines = []
    timeline = []  # [(speaker, index), ...]
    
    lines = [line.strip() for line in script.split('\n') if line.strip()]
    
    for line in lines:
        if ':' in line:
            speaker_part = line.split(':', 1)[0].lower()
            text = line.split(':', 1)[1].strip()
            
            if any(keyword in speaker_part for keyword in ['alex', 'host']):
                host_lines.append(text)
                timeline.append(('host', len(host_lines) - 1))
            elif any(keyword in speaker_part for keyword in ['dr.', 'smith', 'expert', 'researcher']):
                expert_lines.append(text)
                timeline.append(('expert', len(expert_lines) - 1))
    
    return host_lines, expert_lines, timeline


def generate_voice_audio(lines, voice_type, temp_dir, rate):
    """Generate audio for a specific voice type."""
    if not lines:
        return None
    
    output_file = os.path.join(temp_dir, f"{voice_type}_audio.wav")
    combined_text = ' ... '.join(lines)
    
    engine = pyttsx3.init()
    voices = engine.getProperty('voices')
    
    if voice_type == "host" and len(voices) > 1:
        engine.setProperty('voice', voices[1].id)
        engine.setProperty('rate', rate)
    else:
        engine.setProperty('voice', voices[0].id)
        engine.setProperty('rate', rate - 15)
    
    engine.save_to_file(combined_text, output_file)
    engine.runAndWait()
    
    return output_file if os.path.exists(output_file) else None


def interleave_audio_by_timeline(timeline, host_audio, expert_audio, output_file):
    """Interleave audio segments based on timeline."""
    # This is a simplified version - in practice you'd need more sophisticated audio processing
    # For now, just combine the existing files
    if host_audio and expert_audio:
        combine_audio_files([host_audio, expert_audio], output_file)
    elif host_audio:
        import shutil
        shutil.copy2(host_audio, output_file)
    elif expert_audio:
        import shutil
        shutil.copy2(expert_audio, output_file)


def cleanup_temp_files(file_list):
    """Clean up temporary files."""
    for file_path in file_list:
        if file_path and os.path.exists(file_path):
            try:
                os.remove(file_path)
            except:
                pass
                    
@app.route('/process-input', methods=['POST'])
def handle_process_input():
    try:
        data = request.json
        input_type = data.get('input_type')
        input_content = data.get('input_content')
        summary_level = data.get('summary_level', 'Student')
        creativity_level = data.get('creativity_level', 'Balanced')
        podcast_length = data.get('podcast_length', 'Medium (5-7 mins)')
        template_name = data.get('template_name', 'Template 1')

        if not input_type or not input_content:
            return jsonify({'error': 'input_type and input_content are required'}), 400

        # Process based on input type
        if input_type == "PDF":
            # For PDF processing, the frontend should first upload the PDF via /upload-pdf
            # and then send the extracted text here
            text_content = input_content
        else:  # Text input
            # generated_content = generate_content_from_heading(input_content)
            text_content = input_content

        # Generate summary
        summary_text = extract_and_summarize_sections(text_content, summary_level)
        
        # Generate podcast script
        podcast_script = generate_podcast_script(summary_text, creativity_level, podcast_length)
        
        # Generate PPT
        template_path = template_options.get(template_name)
        if not template_path:
            return jsonify({'error': 'Invalid template name'}), 400
        pptx_stream = create_ppt_from_summary(summary_text, template_path)
        
        # Prepare response
        response_data = {
            'summary': summary_text,
            'podcast_script': podcast_script,
            'ppt_file': pptx_stream.getvalue().decode('latin1')  # Encoding for JSON
        }
        
        return jsonify(response_data), 200
        
    except Exception as e:
        return jsonify({
            'error': f'Internal server error: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500

@app.route('/generate-summary', methods=['POST'])
def generate_summary():
    try:
        data = request.json
        text = data.get('text')
        summary_level = data.get('summary_level')

        # Validate input
        if not text or not summary_level:
            return jsonify({'error': 'Invalid input. "text" and "summary_level" are required.'}), 400

        # Call your existing summary generation function
        summary_text = extract_and_summarize_sections(text, summary_level)

        if not summary_text:
            return jsonify({'error': 'Failed to generate summary.'}), 500

        # Split the summary into an array of lines
        summary_lines = [line for line in summary_text.split('\n') if line.strip()]
        
        return jsonify({
            'summary': summary_lines,
            'summary_text': summary_text  # Also return the full text for other uses
        }), 200
    except Exception as e:
        return jsonify({'error': f'Internal server error: {str(e)}'}), 500    
@app.route('/generate-podcast', methods=['POST'])
def generate_podcast():
    try:
        data = request.json
        summary_text = data.get('summary_text')
        creativity_level = data.get('creativity_level')
        podcast_length = data.get('podcast_length')

        if not summary_text or not creativity_level or not podcast_length:
            return jsonify({'error': 'Missing required parameters'}), 400

        # Map frontend values to backend constants
        creativity_map = {
            'formal': 'Formal',
            'balanced': 'Balanced',
            'creative': 'Creative'
        }
        
        length_map = {
            'short': 'Short (2-3 mins)',
            'medium': 'Medium (5-7 mins)',
            'long': 'Long (10+ mins)'
        }

        try:
            mapped_creativity = creativity_map[creativity_level.lower()]
            mapped_length = length_map[podcast_length.lower()]
        except KeyError:
            return jsonify({'error': 'Invalid parameter values'}), 400
        
        # Call your existing podcast generation function
        podcast_script = generate_podcast_script(summary_text, mapped_creativity, mapped_length)
        
        return jsonify({
            'podcast_script': podcast_script,
            'status': 'success'
        }), 200
    except Exception as e:
        return jsonify({
            'error': f'Error generating podcast: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500
                
@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    try:
        data = request.json
        summary_text = data.get('summary_text')
        template_name = data.get('template_name')

        if not summary_text or not template_name:
            return jsonify({'error': 'Missing required parameters'}), 400

        template_path = template_options.get(template_name)
        if not template_path or not os.path.exists(template_path):
            return jsonify({'error': 'Invalid template name'}), 400

        pptx_stream = create_ppt_from_summary(summary_text, template_path)
        
        # Return as binary file
        return send_file(
            pptx_stream,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentation.pptx'
        )
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    
@app.route('/generate-audio', methods=['POST'])
def handle_generate_audio():
    try:
        data = request.json
        podcast_script = data.get('podcast_script')

        if not podcast_script:
            return jsonify({'error': 'podcast_script is required'}), 400

        print(f"Generating audio for podcast script: {podcast_script[:100]}...")  # Log the first 100 characters
        try:
            # Generate audio file
            audio_file = generate_podcast_audio(podcast_script)
            
            # Check if file was created
            if not os.path.exists(audio_file):
                return jsonify({'error': 'Audio file generation failed'}), 500

            # Create a BytesIO buffer to send the file
            with open(audio_file, 'rb') as f:
                audio_bytes = BytesIO(f.read())
            
            # Clean up the temporary file
            os.remove(audio_file)
            
            # Return the audio file
            return send_file(
                audio_bytes,
                mimetype="audio/wav",
                as_attachment=True,
                download_name="podcast_audio.wav"
            )
            
        except Exception as e:
            print(f"Error generating audio: {e}")
            return jsonify({
                'error': f'Audio generation error: {str(e)}',
                'traceback': traceback.format_exc()
            }), 500
            
    except Exception as e:
        print(f"Error in /generate-audio: {e}")
        return jsonify({
            'error': f'Internal server error: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500
                
@app.route('/')
def index():
    return jsonify({'message': 'Backend is running'})

@app.route('/upload-pdf', methods=['POST'])
def handle_upload_pdf():
    try:
        if 'pdf' not in request.files:
            return jsonify({'error': 'No file part in the request'}), 400

        file = request.files['pdf']
        if file.filename == '':
            return jsonify({'error': 'No selected file'}), 400

        if not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': 'Invalid file type. Only PDFs are allowed.'}), 400

        temp_dir = tempfile.mkdtemp()
        temp_pdf_path = os.path.join(temp_dir, secure_filename(file.filename))
        
        try:
            file.save(temp_pdf_path)
            if not os.path.exists(temp_pdf_path):
                return jsonify({'error': 'Failed to save uploaded file'}), 500

            loader = PyPDFLoader(temp_pdf_path)
            documents = loader.load()
            pdf_text = "\n".join([doc.page_content for doc in documents])
            
            return jsonify({
                'message': 'PDF uploaded successfully',
                'content': pdf_text,
                'filename': file.filename
            }), 200
        finally:
            if os.path.exists(temp_pdf_path):
                try:
                    os.remove(temp_pdf_path)
                except:
                    pass
    except Exception as e:
        return jsonify({
            'error': f'Internal server error: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500

# Add these imports at the top of ex.py
import torch
from diffusers import StableDiffusionPipeline
import textwrap
import math
from PIL import Image, ImageDraw, ImageFont
from gradio_client import Client, handle_file

# Add this configuration after other constants
comic_config = {
    "default_prompt": "Two professionals in deep discussion at a sleek, high-tech conference table. Holographic displays show AI models, real-time data, and futuristic UI. Neon lighting reflects off glass surfaces, creating a dynamic, focused atmosphere that highlights AI-driven decision-making.",
    "panel_size": (512, 512),
    "bubble_padding": 40,
    "text_width": 20,
    "font_size": 48
}

# Initialize the Gradio client (add this after other initializations)
# comic_client = Client("gabrielchua/open-notebooklm")

# Initialize Stable Diffusion (add this after other initializations)
# try:
#     comic_pipe = StableDiffusionPipeline.from_pretrained(
#         "stabilityai/stable-diffusion-2", 
#         torch_dtype=torch.float32
#     )
#     comic_pipe.to("cpu")  # Change to 'cuda' if you have GPU support
# except Exception as e:
#     print(f"Warning: Could not load Stable Diffusion model: {str(e)}")
#     comic_pipe = None

# Add these helper functions to ex.py
def get_best_font(size=48):
    """Select the best available font for comic text"""
    possible_fonts = ["ComicSansMS-Bold.ttf", "comicbd.ttf", "arialbd.ttf"]
    for font in possible_fonts:
        try:
            return ImageFont.truetype(font, size)
        except:
            continue
    return ImageFont.load_default()

def generate_comic_image(prompt):
    """Generate a single comic panel image using Stable Diffusion"""
    if not comic_pipe:
        raise Exception("Stable Diffusion model not loaded")
    
    image = comic_pipe(prompt).images[0]
    return image.resize(comic_config["panel_size"])

def add_speech_bubble(img, text, position):
    """Add speech bubble with text to an image"""
    draw = ImageDraw.Draw(img)
    font = get_best_font(size=comic_config["font_size"])
    wrapped_text = textwrap.fill(text, width=comic_config["text_width"])
    
    # Calculate bubble size
    text_bbox = draw.textbbox((0, 0), wrapped_text, font=font)
    bubble_width = text_bbox[2] - text_bbox[0] + comic_config["bubble_padding"]
    bubble_height = text_bbox[3] - text_bbox[1] + comic_config["bubble_padding"]
    
    # Draw bubble
    x, y = position
    draw.rounded_rectangle(
        (x, y, x + bubble_width, y + bubble_height),
        radius=20,
        fill=(255, 255, 255),
        outline="black",
        width=4
    )
    
    # Add text
    draw.text((x + 20, y + 20), wrapped_text, font=font, fill=(0, 0, 0))
    return img

def process_podcast_transcript(file_path):
    """Process podcast audio file to extract dialogues"""
    try:
        result = comic_client.predict(
            files=[handle_file(file_path)],
            url="",
            question="",
            tone="Fun",
            length="Medium (3-5 min)",
            language="English",
            use_advanced_audio=True,
            api_name="/generate_podcast"
        )
        audio_path, transcript_text = result
        dialogues = []
        
        for line in transcript_text.split("\n\n"):
            if ":" in line:
                speaker, text = line.split(": ", 1)
                dialogues.append((speaker.strip(), text.strip()))
        
        return dialogues
    except Exception as e:
        raise Exception(f"Podcast processing failed: {str(e)}")

def generate_comic_from_dialogues(dialogues):
    """Generate comic strip from list of dialogues"""
    if not comic_pipe:
        raise Exception("Stable Diffusion model not loaded")
    
    comic_panels = []
    for char, dialogue in dialogues:
        img = generate_comic_image(comic_config["default_prompt"])
        img = add_speech_bubble(img, dialogue, (20, 20))
        comic_panels.append(img)
    
    # Create comic grid
    num_panels = len(comic_panels)
    cols = min(math.ceil(math.sqrt(num_panels)), 4)
    rows = math.ceil(num_panels / cols)
    
    grid_width = cols * comic_config["panel_size"][0]
    grid_height = rows * comic_config["panel_size"][1]
    comic_grid = Image.new("RGB", (grid_width, grid_height), "white")
    
    for i, panel in enumerate(comic_panels):
        x_offset = (i % cols) * comic_config["panel_size"][0]
        y_offset = (i // cols) * comic_config["panel_size"][1]
        comic_grid.paste(panel, (x_offset, y_offset))
    
    # Save to temporary file
    temp_dir = tempfile.mkdtemp()
    comic_path = os.path.join(temp_dir, "comic.png")
    comic_grid.save(comic_path)
    
    return comic_path

##code snippet got from the claude
@app.route('/generate-comic', methods=['POST'])
def handle_generate_comic():
    try:
        text_content = None
        
        # 1. Extract content from either PDF or text input
        # Check if text content is provided first
        if 'content' in request.form and request.form['content'].strip():
            text_content = request.form['content'].strip()
        elif 'pdf' in request.files:
            # Handle PDF upload
            file = request.files['pdf']
            if file.filename == '':
                return jsonify({'error': 'Empty filename'}), 400
                
            if not file.filename.lower().endswith('.pdf'):
                return jsonify({'error': 'Only PDF files allowed'}), 400

            # Save PDF temporarily and extract text
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
                file.save(temp_pdf.name)
                temp_pdf_path = temp_pdf.name
            
            try:
                loader = PyPDFLoader(temp_pdf_path)
                documents = loader.load()
                text_content = "\n".join([doc.page_content for doc in documents])
                
                if not text_content.strip():
                    return jsonify({'error': 'No readable text found in PDF'}), 400
                    
            except Exception as e:
                return jsonify({'error': f'Failed to extract text from PDF: {str(e)}'}), 400
            finally:
                # Clean up temporary file
                if os.path.exists(temp_pdf_path):
                    try:
                        os.remove(temp_pdf_path)
                    except Exception as e:
                        print(f"Warning: Could not delete temp file: {e}")
        else:
            return jsonify({'error': 'No content provided. Please provide either text content or a PDF file.'}), 400
            
        # Validate we have content
        if not text_content or len(text_content.strip()) < 50:
            return jsonify({'error': 'Insufficient content provided. Please provide at least 50 characters of meaningful text.'}), 400
            
        # 2. Generate natural conversation flow
        dialogue_prompt = f"""
Create a comic-style conversation between two researchers (Alex and Jamie) discussing this research content.
Alex is male and Jamie is female.
Follow this natural flow:
1. Casual greeting and small talk (2-3 exchanges)
2. Transition to research topic ("What are you working on?")
3. Problem statement discussion (2 exchanges)
4. Methodology analysis (3 exchanges)
5. Key findings (2 exchanges)
6. Implications and next steps (2 exchanges)
7. Closing remarks

Rules:
- Keep exchanges 15-25 words each
- Maintain technical accuracy but make it conversational
- Use natural transitions between topics
- Show enthusiasm and curiosity
- Total 10-12 dialogue exchanges
- Format each line as "Speaker: dialogue content"

Research content:
{text_content[:5000]}
"""
        
        try:
            response = model.generate_content(dialogue_prompt)
            dialogue_script = response.text.strip()
            
            # Clean and validate script
            dialogue_lines = []
            for line in dialogue_script.split('\n'):
                line = line.strip()
                if ':' in line and len(line.split(':', 1)) == 2:
                    speaker_part, content = line.split(':', 1)
                    speaker_part = speaker_part.strip()
                    content = content.strip()
                    
                    # Skip empty content
                    if len(content) < 5:
                        continue
                        
                    # Normalize speaker names
                    if "Alex" in speaker_part or "alex" in speaker_part.lower():
                        speaker = "Alex"
                        gender = "male"
                    elif "Jamie" in speaker_part or "jamie" in speaker_part.lower():
                        speaker = "Jamie"
                        gender = "female"
                    else:
                        continue  # Skip unrecognized speakers
                    
                    dialogue_lines.append({
                        "speaker": speaker,
                        "content": content,
                        "gender": gender
                    })
            
            # Ensure we have enough lines - provide fallback dialogue
            if len(dialogue_lines) < 8:
                dialogue_lines = generate_fallback_dialogue(text_content)
                
        except Exception as e:
            print(f"Error generating dialogue: {str(e)}")
            dialogue_lines = generate_fallback_dialogue(text_content)

        # 3. Generate avatar images for each character
        # Pre-define avatar characteristics to maintain consistency
        avatar_characteristics = {
            "Alex": {
                "gender": "male",
                "features": "professional male researcher with short brown hair and glasses, wearing a lab coat",
                "emotion_map": {
                    "greeting": "smiling friendly",
                    "question": "curious expression",
                    "concern": "concerned expression",
                    "excited": "excited expression",
                    "thoughtful": "thoughtful expression"
                }
            },
            "Jamie": {
                "gender": "female",
                "features": "professional female researcher with shoulder-length blonde hair, wearing a blue blouse",
                "emotion_map": {
                    "greeting": "warm smile",
                    "question": "inquisitive expression",
                    "concern": "worried expression",
                    "excited": "enthusiastic expression",
                    "thoughtful": "contemplative expression"
                }
            }
        }

        # Generate comic panels with human conversations and tech backgrounds
        panel_images = []
        base_width, base_height = 512, 512  # Square panels for grid

        # First, generate and cache all avatar images for consistency
        avatar_cache = {}
        for character in ["Alex", "Jamie"]:
            char_info = avatar_characteristics[character]
            for emotion, expression in char_info["emotion_map"].items():
                cache_key = f"{character}_{emotion}"
                
                # Define the avatar prompt for stable diffusion
                avatar_prompt = f"""
                Portrait of a {char_info['features']}, with {expression}, 
                highly detailed realistic face, professional lighting, 
                high quality, photorealistic, 8k, portrait photography
                """
                
                try:
                    if comic_pipe:  # If Stable Diffusion is available
                        avatar_image = comic_pipe(
                            prompt=avatar_prompt,
                            negative_prompt="deformed, ugly, cartoon, anime, blurry, low quality, disfigured",
                            height=base_height,
                            width=base_width,
                            guidance_scale=7.5,
                            num_inference_steps=50
                        ).images[0]
                    else:
                        # Fallback image generation
                        avatar_image = create_placeholder_avatar(character, emotion, base_width, base_height)
                except Exception as e:
                    print(f"Error generating avatar for {cache_key}: {str(e)}")
                    avatar_image = create_placeholder_avatar(character, emotion, base_width, base_height)
                
                avatar_cache[cache_key] = avatar_image

        # Create conversation scene backgrounds
        background_settings = [
            "modern research laboratory with holographic displays",
            "clean white office with large computer monitors",
            "university hallway with research posters",
            "campus coffee shop with laptops and research papers",
            "dimly lit server room with blinking lights",
            "university library with bookshelves and study areas",
            "team meeting room with whiteboard full of diagrams",
            "cybersecurity operations center with multiple screens"
        ]
        
        background_cache = {}
        for i, setting in enumerate(background_settings):
            background_prompt = f"""
            A {setting}, wide angle view, no people, 
            realistic style, professional photography, 
            detailed environment, high quality, 8k
            """
            
            try:
                if comic_pipe:
                    bg_image = comic_pipe(
                        prompt=background_prompt,
                        negative_prompt="low quality, cartoon, anime, people, humans, faces",
                        height=base_height,
                        width=base_width,
                        guidance_scale=7.5,
                        num_inference_steps=40
                    ).images[0]
                else:
                    bg_image = create_placeholder_background(setting, base_width, base_height)
            except Exception as e:
                print(f"Error generating background {i}: {str(e)}")
                bg_image = create_placeholder_background(setting, base_width, base_height)
                
            background_cache[i % len(background_settings)] = bg_image

        # For each dialogue line, create a panel
        for i, line_data in enumerate(dialogue_lines):
            speaker = line_data["speaker"]
            content = line_data["content"]
            gender = line_data["gender"]
            
            # Determine emotion from content
            emotion = determine_emotion(content)
                
            # Get avatar from cache
            avatar_key = f"{speaker}_{emotion}"
            if avatar_key in avatar_cache:
                avatar = avatar_cache[avatar_key]
            else:
                # Fallback to any emotion we have for this character
                for key in avatar_cache:
                    if key.startswith(speaker):
                        avatar = avatar_cache[key]
                        break
                else:
                    avatar = create_placeholder_avatar(speaker, "neutral", base_width, base_height)
            
            # Get background from cache
            background = background_cache[i % len(background_settings)]
            
            # Create panel with both background and avatar
            panel = create_comic_panel(background, avatar, speaker, content, base_width, base_height)
            panel_images.append(panel)

        # 4. Create comic page layout
        comic_page = create_comic_layout(panel_images)

        # Save final comic
        img_io = BytesIO()
        comic_page.save(img_io, 'PNG', quality=95)
        img_io.seek(0)
        
        return send_file(
            img_io,
            mimetype='image/png',
            as_attachment=True,
            download_name='research_comic.png'
        )

    except Exception as e:
        return jsonify({
            'error': f"Failed to generate comic: {str(e)}",
            'traceback': traceback.format_exc()
        }), 500


def generate_fallback_dialogue(text_content):
    """Generate fallback dialogue based on content analysis"""
    # Analyze content for key topics
    content_lower = text_content.lower()
    
    # Determine research domain
    if any(word in content_lower for word in ['security', 'cyber', 'attack', 'vulnerability']):
        topic = "cybersecurity"
    elif any(word in content_lower for word in ['health', 'medical', 'patient', 'clinical']):
        topic = "healthcare"
    elif any(word in content_lower for word in ['data', 'analysis', 'algorithm', 'machine learning']):
        topic = "data science"
    else:
        topic = "research"
    
    # Generate contextual dialogue
    dialogue_templates = {
        "cybersecurity": [
            {"speaker": "Alex", "content": "Hey Jamie! I've been reviewing the latest security assessment.", "gender": "male"},
            {"speaker": "Jamie", "content": "Great timing! I just finished analyzing the vulnerability patterns.", "gender": "female"},
            {"speaker": "Alex", "content": "What did you find regarding the attack vectors?", "gender": "male"},
            {"speaker": "Jamie", "content": "The data shows significant gaps in encryption protocols.", "gender": "female"},
            {"speaker": "Alex", "content": "That's concerning. Are there specific systems at risk?", "gender": "male"},
            {"speaker": "Jamie", "content": "Yes, particularly the legacy systems with outdated security measures.", "gender": "female"},
            {"speaker": "Alex", "content": "We should prioritize updating those systems immediately.", "gender": "male"},
            {"speaker": "Jamie", "content": "Agreed. I'll draft the security recommendations for implementation.", "gender": "female"},
            {"speaker": "Alex", "content": "Perfect. This research could prevent serious breaches.", "gender": "male"},
            {"speaker": "Jamie", "content": "Absolutely. Prevention is always better than response.", "gender": "female"}
        ],
        "healthcare": [
            {"speaker": "Alex", "content": "Jamie, I've been studying the patient data trends.", "gender": "male"},
            {"speaker": "Jamie", "content": "Excellent! I'm analyzing the treatment outcomes as well.", "gender": "female"},
            {"speaker": "Alex", "content": "What patterns are emerging from your analysis?", "gender": "male"},
            {"speaker": "Jamie", "content": "The data suggests improved outcomes with early intervention.", "gender": "female"},
            {"speaker": "Alex", "content": "That aligns with our hypothesis about preventive care.", "gender": "male"},
            {"speaker": "Jamie", "content": "Yes, and the cost-benefit analysis is quite compelling.", "gender": "female"},
            {"speaker": "Alex", "content": "This could significantly impact healthcare policy decisions.", "gender": "male"},
            {"speaker": "Jamie", "content": "We should present these findings to the medical board.", "gender": "female"},
            {"speaker": "Alex", "content": "Great idea. Patient outcomes are our top priority.", "gender": "male"},
            {"speaker": "Jamie", "content": "This research could help save lives and reduce costs.", "gender": "female"}
        ],
        "data science": [
            {"speaker": "Alex", "content": "Jamie, the data analysis results are fascinating!", "gender": "male"},
            {"speaker": "Jamie", "content": "I know! The machine learning models are performing well.", "gender": "female"},
            {"speaker": "Alex", "content": "What's the accuracy rate on the latest algorithm?", "gender": "male"},
            {"speaker": "Jamie", "content": "We're seeing 94% accuracy with the new feature set.", "gender": "female"},
            {"speaker": "Alex", "content": "That's impressive! How did you optimize the parameters?", "gender": "male"},
            {"speaker": "Jamie", "content": "Cross-validation and careful hyperparameter tuning were key.", "gender": "female"},
            {"speaker": "Alex", "content": "The practical applications for this are enormous.", "gender": "male"},
            {"speaker": "Jamie", "content": "Yes, it could revolutionize how we process this data.", "gender": "female"},
            {"speaker": "Alex", "content": "Let's prepare a comprehensive report on our findings.", "gender": "male"},
            {"speaker": "Jamie", "content": "Agreed. This breakthrough deserves proper documentation.", "gender": "female"}
        ],
        "research": [
            {"speaker": "Alex", "content": "Hi Jamie! How's your research project progressing?", "gender": "male"},
            {"speaker": "Jamie", "content": "Really well! I'm seeing some interesting patterns emerge.", "gender": "female"},
            {"speaker": "Alex", "content": "That's exciting! What methodology are you using?", "gender": "male"},
            {"speaker": "Jamie", "content": "A mixed-methods approach with quantitative and qualitative analysis.", "gender": "female"},
            {"speaker": "Alex", "content": "Smart choice. What are your preliminary findings?", "gender": "male"},
            {"speaker": "Jamie", "content": "The data suggests our initial hypothesis was correct.", "gender": "female"},
            {"speaker": "Alex", "content": "Excellent! What are the implications for future work?", "gender": "male"},
            {"speaker": "Jamie", "content": "This opens up several new research avenues to explore.", "gender": "female"},
            {"speaker": "Alex", "content": "We should document these findings thoroughly.", "gender": "male"},
            {"speaker": "Jamie", "content": "Absolutely. This could be the foundation for future studies.", "gender": "female"}
        ]
    }
    
    return dialogue_templates.get(topic, dialogue_templates["research"])


def determine_emotion(content):
    """Determine emotion based on dialogue content"""
    content_lower = content.lower()
    
    if "?" in content:
        return "question"
    elif any(word in content_lower for word in ["concern", "problem", "issue", "worry", "risk"]):
        return "concern"
    elif any(word in content_lower for word in ["great", "excellent", "amazing", "exciting", "perfect"]):
        return "excited"
    elif any(word in content_lower for word in ["hi", "hello", "hey", "good"]):
        return "greeting"
    else:
        return "thoughtful"

def create_comic_panel(background, avatar, speaker, content, base_width, base_height):
    """Create a single comic panel with background, avatar, and speech bubble"""
    panel = Image.new('RGB', (base_width, base_height), (255, 255, 255))
    
    # Paste background first
    panel.paste(background, (0, 0))
    
    # Crop avatar to focus on face and upper body (top 60%)
    avatar_crop = avatar.crop((0, 0, avatar.width, int(avatar.height * 0.6)))
    
    # Resize avatar to fit in panel
    avatar_width = int(base_width * 0.6)
    avatar_height = int(avatar_width * avatar_crop.height / avatar_crop.width)
    avatar_resized = avatar_crop.resize((avatar_width, avatar_height), Image.LANCZOS)
    
    # Position avatar at bottom of panel
    avatar_x = int((base_width - avatar_width) / 2)
    avatar_y = base_height - avatar_height
    
    # Create avatar mask for smooth blending
    mask = Image.new('L', (avatar_width, avatar_height), 0)
    mask_draw = ImageDraw.Draw(mask)
    mask_draw.ellipse((0, 0, avatar_width, avatar_height * 2), fill=255)
    
    # Paste avatar onto panel with mask
    panel.paste(avatar_resized, (avatar_x, avatar_y), mask)
    
    # Draw speech bubble
    draw = ImageDraw.Draw(panel)
    try:
        font = ImageFont.truetype("arial.ttf", 18)
        title_font = ImageFont.truetype("arial.ttf", 22)
    except:
        font = ImageFont.load_default()
        title_font = ImageFont.load_default()

    # Speech bubble at top
    bubble_margin = 20
    bubble_width = base_width - 2 * bubble_margin
    bubble_height = 120
    bubble_y = bubble_margin

    # Draw bubble with color based on speaker
    bubble_color = "#3377cc" if speaker == "Alex" else "#cc3377"
    
    draw.rounded_rectangle(
        [bubble_margin, bubble_y, bubble_margin + bubble_width, bubble_y + bubble_height],
        radius=15, fill="white", outline=bubble_color, width=3
    )

    # Bubble tail pointing to speaker
    tail_x = base_width // 2
    draw.polygon([
        (tail_x - 15, bubble_y + bubble_height),
        (tail_x + 15, bubble_y + bubble_height),
        (tail_x, bubble_y + bubble_height + 15)
    ], fill="white", outline=bubble_color)

    # Speaker name in color
    draw.text((bubble_margin + 15, bubble_y + 10), speaker, fill=bubble_color, font=title_font)

    # Content text
    wrapped_text = textwrap.fill(content, width=32)
    draw.text((bubble_margin + 15, bubble_y + 40), wrapped_text, fill="black", font=font)

    return panel


def create_comic_layout(panel_images):
    """Create the final comic page layout"""
    page_width = 1100
    page_margins = 50
    panel_margin = 20
    panels_per_row = 2
    panel_width = (page_width - 2 * page_margins - (panels_per_row - 1) * panel_margin) // panels_per_row
    
    # Calculate height needed for each panel (maintaining aspect ratio)
    panel_height = panel_width
    
    # Calculate total height needed
    rows = math.ceil(len(panel_images) / panels_per_row)
    page_height = 2 * page_margins + rows * panel_height + (rows - 1) * panel_margin + 60  # Extra space for title
    
    # Create the comic page
    comic_page = Image.new('RGB', (page_width, page_height), (240, 240, 240))
    
    # Add title at the top
    draw = ImageDraw.Draw(comic_page)
    try:
        title_font = ImageFont.truetype("arial.ttf", 36)
    except:
        title_font = ImageFont.load_default()
        
    title = "Research Discussion Comic"
    # For compatibility, calculate title dimensions manually
    title_width = len(title) * 20  # Approximate width
    title_height = 40
    draw.text(((page_width - title_width) // 2, page_margins // 2), title, fill=(0, 0, 0), font=title_font)
    
    # Place each panel on the page
    for i, panel in enumerate(panel_images):
        row = i // panels_per_row
        col = i % panels_per_row
        
        x = page_margins + col * (panel_width + panel_margin)
        y = page_margins + 60 + row * (panel_height + panel_margin)  # 60px offset for title
        
        # Resize panel to fit layout
        resized_panel = panel.resize((panel_width, panel_height), Image.LANCZOS)
        comic_page.paste(resized_panel, (x, y))
        
        # Add panel border
        draw.rectangle([x, y, x + panel_width, y + panel_height], outline=(0, 0, 0), width=2)
    
    return comic_page


def create_placeholder_avatar(character, emotion, width, height):
    """Create a placeholder avatar when AI generation fails"""
    img = Image.new('RGB', (width, height), (200, 200, 200))
    draw = ImageDraw.Draw(img)
    
    try:
        font = ImageFont.truetype("arial.ttf", 24)
    except:
        font = ImageFont.load_default()
    
    text = f"{character}\n({emotion})"
    draw.text((width//4, height//2), text, fill=(0, 0, 0), font=font)
    
    return img


def create_placeholder_background(setting, width, height):
    """Create a placeholder background when AI generation fails"""
    img = Image.new('RGB', (width, height), (150, 150, 150))
    draw = ImageDraw.Draw(img)
    
    try:
        font = ImageFont.truetype("arial.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    wrapped_setting = textwrap.fill(setting, width=20)
    draw.text((20, height//2), wrapped_setting, fill=(255, 255, 255), font=font)
    
    return img

def create_placeholder_avatar(character, emotion, width, height):
    """Create a placeholder avatar when image generation fails"""
    avatar = Image.new('RGB', (width, height), (200, 200, 220))
    draw = ImageDraw.Draw(avatar)
    
    # Draw face outline
    face_color = (255, 220, 200)
    draw.ellipse([width//4, height//4, width*3//4, height*3//4], fill=face_color)
    
    # Draw eyes
    eye_color = (50, 100, 150) if character == "Alex" else (100, 50, 150)
    draw.ellipse([width//3, height*2//5, width*2//5, height//2], fill=(255, 255, 255), outline=(0, 0, 0))
    draw.ellipse([width*3//5, height*2//5, width*2//3, height//2], fill=(255, 255, 255), outline=(0, 0, 0))
    draw.ellipse([width//3+5, height*2//5+5, width*2//5-5, height//2-5], fill=eye_color)
    draw.ellipse([width*3//5+5, height*2//5+5, width*2//3-5, height//2-5], fill=eye_color)
    
    # Draw mouth (vary by emotion)
    if emotion in ["greeting", "excited"]:
        # Smile
        draw.arc([width//3, height//2, width*2//3, height*2//3], 0, 180, fill=(0, 0, 0), width=3)
    elif emotion == "concern":
        # Frown
        draw.arc([width//3, height*7//12, width*2//3, height*5//6], 180, 360, fill=(0, 0, 0), width=3)
    elif emotion == "question":
        # Slightly open mouth
        draw.ellipse([width*2//5, height*3//5, width*3//5, height*2//3], fill=(150, 50, 50), outline=(0, 0, 0))
    else:
        # Neutral
        draw.line([width//3, height*3//5, width*2//3, height*3//5], fill=(0, 0, 0), width=3)
    
    # Add hair
    hair_color = (100, 70, 40) if character == "Alex" else (220, 190, 140)
    if character == "Alex":  # Short hair
        draw.rectangle([width//4, height//5, width*3//4, height*2//5], fill=hair_color)
        # Draw glasses
        draw.rectangle([width//3-5, height*2//5, width*2//5+5, height//2], outline=(0, 0, 0), width=2)
        draw.rectangle([width*3//5-5, height*2//5, width*2//3+5, height//2], outline=(0, 0, 0), width=2)
        draw.line([width*2//5+5, height*9//20, width*3//5-5, height*9//20], fill=(0, 0, 0), width=2)
    else:  # Jamie - longer hair
        draw.rectangle([width//5, height//5, width*4//5, height*2//5], fill=hair_color)
        draw.rectangle([width//6, height*2//5, width//4, height*3//5], fill=hair_color)
        draw.rectangle([width*3//4, height*2//5, width*5//6, height*3//5], fill=hair_color)
    
    # Add character name
    try:
        font = ImageFont.truetype("arial.ttf", 24)
    except:
        font = ImageFont.load_default()
    draw.text((width//2, height*4//5), character, fill=(0, 0, 0), font=font, anchor="mm")
    
    return avatar

def create_placeholder_background(setting, width, height):
    """Create a placeholder background when image generation fails"""
    colors = {
        "laboratory": (220, 230, 250),
        "office": (240, 240, 240),
        "hallway": (230, 225, 200),
        "coffee": (210, 190, 170),
        "server": (50, 60, 80),
        "library": (210, 200, 180),
        "meeting": (230, 230, 240),
        "center": (60, 70, 90)
    }
    
    # Determine color based on setting keyword
    bg_color = (220, 220, 220)  # default
    for keyword, color in colors.items():
        if keyword in setting.lower():
            bg_color = color
            break
    
    background = Image.new('RGB', (width, height), bg_color)
    draw = ImageDraw.Draw(background)
    
    # Add simple elements based on setting
    if "laboratory" in setting or "research" in setting:
        # Draw lab equipment outlines
        for i in range(5):
            x = width // 5 * i + width // 10
            y = height // 3
            draw.rectangle([x-30, y-30, x+30, y+30], outline=(100, 100, 100), width=2)
            draw.line([x, y-30, x, y+30], fill=(100, 100, 100), width=2)
    
    elif "office" in setting or "computer" in setting:
        # Draw computer monitors
        for i in range(3):
            x = width // 4 * (i+1)
            y = height // 3
            draw.rectangle([x-40, y-30, x+40, y+30], outline=(80, 80, 80), width=3)
            draw.rectangle([x-10, y+30, x+10, y+40], outline=(80, 80, 80), width=2)
            draw.rectangle([x-20, y+40, x+20, y+45], fill=(80, 80, 80))
    
    elif "server" in setting:
        # Draw server racks
        for i in range(5):
            x = width // 6 * (i+1)
            draw.rectangle([x-30, height//5, x+30, height*4//5], fill=(40, 40, 40), outline=(30, 30, 30))
            for j in range(10):
                y = height//5 + j * height//20
                draw.line([x-25, y, x+25, y], fill=(60, 60, 60), width=1)
                draw.rectangle([x+15, y+2, x+20, y+5], fill=(0, 255, 0))  # Green LED
    
    else:
        # Generic background with grid
        for i in range(0, width, 40):
            draw.line([i, 0, i, height], fill=(max(bg_color[0]-30, 0), max(bg_color[1]-30, 0), max(bg_color[2]-30, 0)), width=1)
        for i in range(0, height, 40):
            draw.line([0, i, width, i], fill=(max(bg_color[0]-30, 0), max(bg_color[1]-30, 0), max(bg_color[2]-30, 0)), width=1)
    
    # Add text label at bottom
    try:
        font = ImageFont.truetype("arial.ttf", 16)
    except:
        font = ImageFont.load_default()
    
    setting_short = setting[:20] + "..." if len(setting) > 20 else setting
    draw.rectangle([0, height-25, width, height], fill=(0, 0, 0, 128))
    draw.text((width//2, height-12), setting_short, fill=(255, 255, 255), font=font, anchor="mm")
    
    return background

## a bit good approach for video
import os
import re
import tempfile
import traceback
import subprocess
import random
import math
from flask import request, jsonify, send_file
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import textwrap

class AvatarAnimator:
    def __init__(self, width, height):
        self.width = width
        self.height = height
        self.avatar_size = 120
        self.avatar_colors = [
            (255, 192, 203),  # Pink
            (135, 206, 235),  # Sky Blue
            (255, 165, 0),    # Orange
            (144, 238, 144),  # Light Green
            (221, 160, 221),  # Plum
            (255, 215, 0),    # Gold
        ]
        
    def create_avatar(self, color, expression="happy", frame=0):
        """Create an animated avatar with different expressions"""
        avatar = Image.new('RGBA', (self.avatar_size, self.avatar_size), (0, 0, 0, 0))
        draw = ImageDraw.Draw(avatar)
        
        # Body (circle)
        body_offset = int(5 * math.sin(frame * 0.1))  # Gentle floating animation
        draw.ellipse([10, 20 + body_offset, self.avatar_size-10, self.avatar_size-20 + body_offset], 
                    fill=color, outline=(0, 0, 0, 100), width=2)
        
        # Eyes
        eye_y = 40 + body_offset
        if expression == "happy":
            # Happy eyes (crescents)
            draw.arc([25, eye_y, 40, eye_y+15], 0, 180, fill=(0, 0, 0), width=3)
            draw.arc([70, eye_y, 85, eye_y+15], 0, 180, fill=(0, 0, 0), width=3)
        elif expression == "excited":
            # Wide open eyes
            draw.ellipse([25, eye_y, 40, eye_y+15], fill=(255, 255, 255), outline=(0, 0, 0), width=2)
            draw.ellipse([70, eye_y, 85, eye_y+15], fill=(255, 255, 255), outline=(0, 0, 0), width=2)
            draw.ellipse([30, eye_y+3, 35, eye_y+12], fill=(0, 0, 0))
            draw.ellipse([75, eye_y+3, 80, eye_y+12], fill=(0, 0, 0))
        else:  # normal
            draw.ellipse([25, eye_y, 40, eye_y+15], fill=(255, 255, 255), outline=(0, 0, 0), width=2)
            draw.ellipse([70, eye_y, 85, eye_y+15], fill=(255, 255, 255), outline=(0, 0, 0), width=2)
            draw.ellipse([28, eye_y+5, 37, eye_y+10], fill=(0, 0, 0))
            draw.ellipse([73, eye_y+5, 82, eye_y+10], fill=(0, 0, 0))
        
        # Mouth
        mouth_y = 70 + body_offset
        if expression == "happy" or expression == "excited":
            # Smile
            draw.arc([40, mouth_y, 70, mouth_y+20], 0, 180, fill=(0, 0, 0), width=3)
        else:
            # Neutral mouth
            draw.ellipse([50, mouth_y+5, 60, mouth_y+10], fill=(0, 0, 0))
        
        # Arms (animated waving)
        arm_angle = 20 * math.sin(frame * 0.2)
        left_arm_x = 15 + int(10 * math.sin(math.radians(arm_angle)))
        right_arm_x = 95 - int(10 * math.sin(math.radians(arm_angle)))
        
        draw.ellipse([left_arm_x, 35 + body_offset, left_arm_x+15, 50 + body_offset], 
                    fill=color, outline=(0, 0, 0), width=1)
        draw.ellipse([right_arm_x, 35 + body_offset, right_arm_x+15, 50 + body_offset], 
                    fill=color, outline=(0, 0, 0), width=1)
        
        return avatar
    
    def create_speech_bubble(self, text, position, frame=0):
        """Create an animated speech bubble"""
        # Wrap text
        wrapper = textwrap.TextWrapper(width=20)
        lines = wrapper.wrap(text)[:3]  # Max 3 lines
        
        # Calculate bubble size
        bubble_width = max(200, max(len(line) * 8 for line in lines) + 40)
        bubble_height = len(lines) * 25 + 30
        
        # Create bubble with animation (slight scaling)
        scale = 1 + 0.05 * math.sin(frame * 0.15)
        scaled_width = int(bubble_width * scale)
        scaled_height = int(bubble_height * scale)
        
        bubble = Image.new('RGBA', (scaled_width + 30, scaled_height + 30), (0, 0, 0, 0))
        draw = ImageDraw.Draw(bubble)
        
        # Bubble background
        draw.rounded_rectangle([15, 5, scaled_width, scaled_height], 
                             radius=15, fill=(255, 255, 255, 240), 
                             outline=(0, 0, 0, 180), width=2)
        
        # Bubble tail
        draw.polygon([(30, scaled_height), (45, scaled_height), (35, scaled_height + 15)], 
                    fill=(255, 255, 255, 240), outline=(0, 0, 0, 180))
        
        # Text
        try:
            font = ImageFont.truetype("arial.ttf", 16)
        except IOError:
            font = ImageFont.load_default()
        
        y_offset = 15
        for line in lines:
            text_width = draw.textlength(line, font=font) if hasattr(draw, 'textlength') else len(line) * 8
            x_offset = (scaled_width - text_width) // 2 + 15
            draw.text((x_offset, y_offset), line, fill=(0, 0, 0), font=font)
            y_offset += 25
        
        return bubble

class EnhancedVideoGenerator:
    def __init__(self):
        self.animator = None
        
    def analyze_content_mood(self, text):
        """Analyze the content to determine the appropriate mood and style"""
        excitement_words = ['amazing', 'incredible', 'awesome', 'fantastic', 'exciting', 'wow', 'great']
        educational_words = ['learn', 'understand', 'study', 'knowledge', 'science', 'research']
        funny_words = ['funny', 'hilarious', 'joke', 'laugh', 'humor', 'comedy']
        
        text_lower = text.lower()
        
        if any(word in text_lower for word in excitement_words):
            return 'excited'
        elif any(word in text_lower for word in funny_words):
            return 'funny'
        elif any(word in text_lower for word in educational_words):
            return 'educational'
        else:
            return 'normal'
    
    def add_background_elements(self, img, style, frame):
        """Add animated background elements"""
        draw = ImageDraw.Draw(img)
        width, height = img.size
        
        if style == 'modern':
            # Floating geometric shapes
            for i in range(5):
                x = (width // 6) * i + int(20 * math.sin(frame * 0.05 + i))
                y = height - 100 + int(15 * math.cos(frame * 0.08 + i))
                size = 20 + int(10 * math.sin(frame * 0.1 + i))
                
                # Create semi-transparent overlay
                overlay = Image.new('RGBA', img.size, (0, 0, 0, 0))
                overlay_draw = ImageDraw.Draw(overlay)
                
                if i % 2 == 0:
                    # Circles
                    overlay_draw.ellipse([x, y, x+size, y+size], 
                                       fill=(255, 255, 255, 50))
                else:
                    # Squares
                    overlay_draw.rectangle([x, y, x+size, y+size], 
                                         fill=(200, 200, 255, 60))
                
                img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
        
        return img
    
    def create_transition_effect(self, img1, img2, progress):
        """Create smooth transitions between slides"""
        # Simple fade transition
        alpha = int(255 * progress)
        img2_faded = img2.copy()
        img2_faded.putalpha(alpha)
        
        result = img1.copy().convert('RGBA')
        result = Image.alpha_composite(result, img2_faded.convert('RGBA'))
        return result.convert('RGB')

@app.route('/generate-video', methods=['POST'])
def generate_video():
    try:
        # Get the summary text from the request
        data = request.get_json()
        if not data or 'summary_text' not in data:
            return jsonify({'error': 'No summary text provided'}), 400
        
        summary_text = data['summary_text']
        
        # Optional parameters with defaults
        video_style = data.get('video_style', 'modern')
        resolution = data.get('resolution', '720p')
        
        # Create enhanced video generator
        generator = EnhancedVideoGenerator()
        
        # Create a temporary directory for the video generation
        temp_dir = tempfile.mkdtemp()
        frames_dir = os.path.join(temp_dir, "frames")
        os.makedirs(frames_dir, exist_ok=True)
        output_path = os.path.join(temp_dir, "summary_video.mp4")
        
        try:
            # Set resolution
            if resolution == '720p':
                width, height = 1280, 720
            elif resolution == '1080p':
                width, height = 1920, 1080
            else:
                width, height = 1280, 720  # Default
            
            generator.animator = AvatarAnimator(width, height)
            
            # Enhanced key point extraction function
            def extract_key_points(text, max_points=5):
                """Extract key points from text using multiple methods"""
                sentences = re.split(r'[.!?]+', text)
                sentences = [s.strip() for s in sentences if len(s.strip()) > 10]
                
                key_points = []
                
                # Method 1: Look for sentences with keywords
                priority_words = ['important', 'key', 'main', 'crucial', 'significant', 
                                'result', 'conclusion', 'because', 'therefore', 'however']
                
                for sentence in sentences:
                    if any(word in sentence.lower() for word in priority_words):
                        key_points.append(sentence.strip())
                
                # Method 2: Look for numbered/bulleted points
                bullet_patterns = [r'^\d+\.', r'^[•\-\*]', r'^[a-z]\)', r'^\([a-z]\)']
                for sentence in sentences:
                    if any(re.match(pattern, sentence.strip()) for pattern in bullet_patterns):
                        key_points.append(re.sub(r'^[\d\.\-\*•\(\)a-z]+\s*', '', sentence.strip()))
                
                # Method 3: Extract sentences with specific structures
                for sentence in sentences:
                    if (len(sentence.split()) >= 5 and len(sentence.split()) <= 25 and
                        ('is' in sentence or 'are' in sentence or 'can' in sentence or 'will' in sentence)):
                        key_points.append(sentence.strip())
                
                # Remove duplicates and filter
                unique_points = []
                for point in key_points:
                    if point not in unique_points and len(point) > 20:
                        unique_points.append(point)
                
                # If we don't have enough points, add the longest sentences
                if len(unique_points) < max_points:
                    remaining_sentences = [s for s in sentences if s not in unique_points]
                    remaining_sentences.sort(key=len, reverse=True)
                    for sentence in remaining_sentences[:max_points - len(unique_points)]:
                        if len(sentence) > 30:
                            unique_points.append(sentence)
                
                return unique_points[:max_points]
            
            # Parse content into slides with key points
            slides = []
            sections = re.split(r'\n\s*##\s+', summary_text)
            
            # Process sections and extract key points
            if not sections[0].strip().startswith('# '):
                intro_text = sections[0].strip()
                if intro_text:
                    intro_points = extract_key_points(intro_text, 3)
                    slides.append({
                        'title': '🎬 Welcome to Our Story!',
                        'key_points': intro_points,
                        'mood': generator.analyze_content_mood(intro_text),
                        'avatar_count': 2,
                        'animation_style': 'fade_in'
                    })
                sections = sections[1:]
            
            # Process remaining sections
            for section in sections:
                lines = section.strip().split('\n')
                if not lines:
                    continue
                    
                title = lines[0].strip()
                content = '\n'.join(lines[1:]).strip()
                full_content = content + ' ' + title
                
                key_points = extract_key_points(full_content, 4)
                animation_styles = ['slide_in', 'bounce_in', 'zoom_in', 'flip_in', 'typewriter']
                
                slides.append({
                    'title': title,
                    'key_points': key_points,
                    'mood': generator.analyze_content_mood(full_content),
                    'avatar_count': random.randint(2, 3),
                    'animation_style': random.choice(animation_styles)
                })
            
            # Try to load fonts
            try:
                title_font = ImageFont.truetype("arial.ttf", 48)
                content_font = ImageFont.truetype("arial.ttf", 24)
                point_font = ImageFont.truetype("arial.ttf", 20)
                small_font = ImageFont.truetype("arial.ttf", 16)
            except IOError:
                title_font = ImageFont.load_default()
                content_font = ImageFont.load_default()
                point_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
            
            # Define enhanced styles
            if video_style == 'modern':
                bg_gradient_start = (45, 45, 85)
                bg_gradient_end = (25, 25, 55)
                title_color = (255, 255, 255)
                point_colors = [(255, 215, 0), (255, 105, 180), (0, 255, 127), (255, 69, 0), (138, 43, 226)]
                content_bg = (240, 245, 255, 200)
            else:
                bg_gradient_start = (20, 50, 100)
                bg_gradient_end = (0, 20, 60)
                title_color = (255, 255, 255)
                point_colors = [(255, 255, 0), (255, 165, 0), (0, 255, 255), (255, 20, 147), (124, 252, 0)]
                content_bg = (255, 255, 255, 220)
            
            frame_count = 0
            
            # Animation helper functions
            def apply_animation(text, animation_style, progress, max_chars):
                """Apply different animation styles to text appearance"""
                if animation_style == 'typewriter':
                    return text[:int(progress * len(text))]
                elif animation_style == 'fade_in':
                    if progress > 0.3:
                        return text
                    else:
                        return ""
                elif animation_style == 'slide_in':
                    if progress > 0.2:
                        return text
                    else:
                        return ""
                elif animation_style == 'bounce_in':
                    if progress > 0.4:
                        return text
                    else:
                        return ""
                elif animation_style == 'zoom_in':
                    if progress > 0.3:
                        return text
                    else:
                        return ""
                elif animation_style == 'flip_in':
                    if progress > 0.5:
                        return text
                    else:
                        return ""
                return text
            
            def get_point_position(point_idx, total_points, frame, width, height, animation_style):
                """Calculate animated positions for key points"""
                base_y = 180 + point_idx * 120
                base_x = 100
                
                if animation_style == 'slide_in':
                    offset_x = int(50 * math.cos(frame * 0.02 + point_idx))
                    return (base_x + offset_x, base_y)
                elif animation_style == 'bounce_in':
                    bounce = int(10 * abs(math.sin(frame * 0.1 + point_idx * 0.5)))
                    return (base_x, base_y - bounce)
                elif animation_style == 'zoom_in':
                    pulse = int(5 * math.sin(frame * 0.08 + point_idx))
                    return (base_x + pulse, base_y + pulse)
                elif animation_style == 'flip_in':
                    flip_offset = int(15 * math.sin(frame * 0.05 + point_idx * 0.8))
                    return (base_x + flip_offset, base_y)
                else:  # typewriter or default
                    return (base_x, base_y)
            
            # Create animated title sequence
            title_duration = 4
            for i in range(24 * title_duration):
                img = Image.new('RGB', (width, height), bg_gradient_start)
                draw = ImageDraw.Draw(img)
                
                # Animated title with particle effects
                progress = min(1.0, i / (24 * 2))
                title_alpha = int(255 * progress)
                
                main_title = "🎬 AI Key Points Video Generator 🎬"
                
                # Create title overlay with glow effect
                title_overlay = Image.new('RGBA', (width, height), (0, 0, 0, 0))
                title_draw = ImageDraw.Draw(title_overlay)
                
                title_bbox = title_draw.textbbox((0, 0), main_title, font=title_font)
                title_w = title_bbox[2] - title_bbox[0]
                title_x = (width - title_w) // 2
                title_y = height // 2 - 50
                
                # Add glow effect
                for offset in range(5):
                    alpha = max(0, title_alpha - offset * 50)
                    title_draw.text((title_x + offset, title_y + offset), main_title, 
                                  fill=(100, 100, 255, alpha // 3), font=title_font)
                
                title_draw.text((title_x, title_y), main_title, 
                              fill=(*title_color, title_alpha), font=title_font)
                
                # Add floating particles
                if progress > 0.5:
                    for j in range(20):
                        particle_x = random.randint(0, width)
                        particle_y = random.randint(0, height)
                        particle_size = random.randint(1, 4)
                        particle_alpha = random.randint(100, 255)
                        color = random.choice(point_colors)
                        title_draw.ellipse([particle_x, particle_y, 
                                          particle_x + particle_size, particle_y + particle_size],
                                         fill=(*color, particle_alpha))
                
                img = Image.alpha_composite(img.convert('RGBA'), title_overlay).convert('RGB')
                img = generator.add_background_elements(img, video_style, i)
                
                img.save(os.path.join(frames_dir, f"frame_{frame_count:05d}.jpg"), quality=95)
                frame_count += 1
            
            # Create slides with animated key points
            for slide_idx, slide in enumerate(slides):
                key_points = slide['key_points']
                slide_duration = max(10, len(key_points) * 4)  # 4 seconds per key point minimum
                
                # Create avatars for this slide
                avatars = []
                avatar_colors = random.sample(generator.animator.avatar_colors, slide['avatar_count'])
                
                for avatar_idx in range(slide['avatar_count']):
                    avatars.append({
                        'color': avatar_colors[avatar_idx],
                        'position': (150 + avatar_idx * 250, height - 180),
                        'expression': slide['mood'] if slide['mood'] in ['happy', 'excited'] else 'happy'
                    })
                
                for frame in range(24 * slide_duration):
                    img = Image.new('RGB', (width, height), bg_gradient_start)
                    img = generator.add_background_elements(img, video_style, frame_count + frame)
                    
                    # Create overlay for animations
                    overlay = Image.new('RGBA', (width, height), (0, 0, 0, 0))
                    overlay_draw = ImageDraw.Draw(overlay)
                    
                    # Draw animated title
                    title_y_offset = int(10 * math.sin((frame_count + frame) * 0.05))
                    overlay_draw.rectangle([(0, 0 + title_y_offset), (width, 80 + title_y_offset)], 
                                         fill=(*bg_gradient_end, 180))
                    
                    # Add emoji based on mood
                    emoji_map = {
                        'excited': '🎉',
                        'funny': '😄',
                        'educational': '📚',
                        'happy': '😊',
                        'default': '💡'
                    }
                    emoji = emoji_map.get(slide['mood'], '💡')
                    
                    full_title = f"{emoji} {slide['title']} {emoji}"
                    title_bbox = overlay_draw.textbbox((0, 0), full_title, font=title_font)
                    title_w = title_bbox[2] - title_bbox[0]
                    overlay_draw.text(((width - title_w) // 2, 20 + title_y_offset), full_title, 
                                    fill=title_color, font=title_font)
                    
                    # Animate key points with different styles
                    for point_idx, point in enumerate(key_points):
                        if len(point.strip()) == 0:
                            continue
                            
                        # Calculate timing for each point
                        point_start_time = point_idx * (slide_duration / len(key_points))
                        point_duration = slide_duration / len(key_points) + 2
                        current_time = frame / 24.0
                        
                        if current_time >= point_start_time:
                            # Calculate progress for this point
                            point_progress = min(1.0, (current_time - point_start_time) / point_duration)
                            
                            # Get animated position
                            pos_x, pos_y = get_point_position(point_idx, len(key_points), frame, 
                                                             width, height, slide['animation_style'])
                            
                            # Apply text animation
                            animated_text = apply_animation(point, slide['animation_style'], 
                                                          point_progress, len(point))
                            
                            if animated_text:
                                # Create point background with animation effects
                                point_color = point_colors[point_idx % len(point_colors)]
                                
                                # Animated background for point
                                bg_alpha = int(200 * point_progress)
                                bg_width = min(width - 200, len(animated_text) * 12)
                                bg_height = 60
                                
                                # Add pulsing effect
                                pulse = int(5 * math.sin(frame * 0.1 + point_idx))
                                
                                # Background with rounded corners effect
                                overlay_draw.rounded_rectangle([
                                    (pos_x - 20 + pulse, pos_y - 15 + pulse),
                                    (pos_x + bg_width + pulse, pos_y + bg_height - 15 + pulse)
                                ], radius=15, fill=(*point_color, bg_alpha // 2))
                                
                                # Add bullet point with animation
                                bullet_symbols = ['●', '▶', '✦', '◆', '▸']
                                bullet = bullet_symbols[point_idx % len(bullet_symbols)]
                                
                                # Bullet point with glow
                                bullet_size = int(20 + 5 * math.sin(frame * 0.08 + point_idx))
                                overlay_draw.text((pos_x, pos_y), bullet, 
                                                fill=point_color, font=point_font)
                                
                                # Wrap text properly
                                words = animated_text.split()
                                lines = []
                                current_line = []
                                
                                for word in words:
                                    test_line = ' '.join(current_line + [word])
                                    if len(test_line) > 60:  # Wrap at 60 characters
                                        if current_line:
                                            lines.append(' '.join(current_line))
                                            current_line = [word]
                                        else:
                                            lines.append(word)
                                    else:
                                        current_line.append(word)
                                
                                if current_line:
                                    lines.append(' '.join(current_line))
                                
                                # Draw text lines with shadow effect
                                for line_idx, line in enumerate(lines[:3]):  # Max 3 lines per point
                                    text_y = pos_y + line_idx * 25
                                    
                                    # Shadow
                                    overlay_draw.text((pos_x + 32, text_y + 2), line, 
                                                    fill=(0, 0, 0, 128), font=point_font)
                                    # Main text
                                    overlay_draw.text((pos_x + 30, text_y), line, 
                                                    fill=(255, 255, 255), font=point_font)
                                
                                # Add sparkle effects for completed points
                                if point_progress > 0.8:
                                    for sparkle in range(3):
                                        sparkle_x = pos_x + random.randint(-10, bg_width + 10)
                                        sparkle_y = pos_y + random.randint(-10, bg_height + 10)
                                        sparkle_size = random.randint(2, 5)
                                        overlay_draw.ellipse([
                                            sparkle_x, sparkle_y, 
                                            sparkle_x + sparkle_size, sparkle_y + sparkle_size
                                        ], fill=(*point_color, random.randint(150, 255)))
                    
                    # Composite the overlay
                    img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
                    draw = ImageDraw.Draw(img)
                    
                    # Add animated avatars
                    for avatar_idx, avatar_info in enumerate(avatars):
                        avatar = generator.animator.create_avatar(
                            avatar_info['color'], 
                            avatar_info['expression'], 
                            frame_count + frame + avatar_idx * 10
                        )
                        
                        # Avatar position with complex animation
                        x, base_y = avatar_info['position']
                        y = base_y + int(15 * math.sin((frame_count + frame + avatar_idx * 20) * 0.1))
                        
                        # Add side-to-side movement
                        x += int(20 * math.cos((frame_count + frame + avatar_idx * 30) * 0.05))
                        
                        # Ensure avatar stays within bounds
                        x = max(0, min(x, width - generator.animator.avatar_size))
                        y = max(0, min(y, height - generator.animator.avatar_size))
                        
                        img.paste(avatar, (x, y), avatar)
                        
                        # Dynamic speech bubbles for key points
                        if frame % 180 == avatar_idx * 60:  # Staggered comments every 7.5 seconds
                            reactions = [
                                "Great point! 💡", "Interesting! 🤔", "I see! 👀", 
                                "Amazing! ✨", "Got it! 👍", "Wow! 🤯",
                                "Makes sense! 🧠", "Brilliant! 🌟", "Exactly! ✅"
                            ]
                            reaction = random.choice(reactions)
                            bubble = generator.animator.create_speech_bubble(
                                reaction, (x, y - 50), frame_count + frame
                            )
                            bubble_x = max(0, min(x - 25, width - bubble.width))
                            bubble_y = max(0, y - bubble.height - 20)
                            img.paste(bubble, (bubble_x, bubble_y), bubble)
                    
                    # Enhanced progress indicator
                    slide_progress = frame / (24 * slide_duration)
                    progress_width = int((width - 200) * slide_progress)
                    
                    # Progress bar background
                    draw.rectangle([(100, height - 35), (width - 100, height - 25)], fill=(60, 60, 60))
                    
                    # Animated progress bar
                    gradient_colors = point_colors[:3]
                    for i in range(progress_width):
                        color_idx = int((i / progress_width) * (len(gradient_colors) - 1)) if progress_width > 0 else 0
                        color = gradient_colors[color_idx]
                        draw.line([(100 + i, height - 35), (100 + i, height - 25)], fill=color)
                    
                    # Progress text
                    progress_text = f"Key Point {min(len(key_points), int(slide_progress * len(key_points)) + 1)} of {len(key_points)}"
                    draw.text((width // 2 - 60, height - 50), progress_text, fill=(255, 255, 255), font=small_font)
                    
                    img.save(os.path.join(frames_dir, f"frame_{frame_count:05d}.jpg"), quality=95)
                    frame_count += 1
            
            # Enhanced ending sequence
            ending_duration = 4
            for i in range(24 * ending_duration):
                img = Image.new('RGB', (width, height), bg_gradient_start)
                img = generator.add_background_elements(img, video_style, frame_count + i)
                
                overlay = Image.new('RGBA', (width, height), (0, 0, 0, 0))
                overlay_draw = ImageDraw.Draw(overlay)
                
                # Animated thank you message
                thank_you = "Thanks for watching! 🎬✨"
                subtitle = "Key points delivered with style! 💫"
                
                title_bbox = overlay_draw.textbbox((0, 0), thank_you, font=title_font)
                title_w = title_bbox[2] - title_bbox[0]
                
                # Main title with rainbow effect
                for j, char in enumerate(thank_you):
                    char_color = point_colors[j % len(point_colors)]
                    char_x = (width - title_w) // 2 + j * (title_w // len(thank_you))
                    char_y = height // 2 - 60 + int(10 * math.sin((frame_count + i + j * 5) * 0.1))
                    overlay_draw.text((char_x, char_y), char, fill=char_color, font=title_font)
                
                # Subtitle
                sub_bbox = overlay_draw.textbbox((0, 0), subtitle, font=content_font)
                sub_w = sub_bbox[2] - sub_bbox[0]
                overlay_draw.text(((width - sub_w) // 2, height // 2 + 20), subtitle, 
                                fill=(255, 255, 255), font=content_font)
                
                img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
                draw = ImageDraw.Draw(img)
                
                # Final celebration avatars
                for j in range(4):
                    avatar = generator.animator.create_avatar(
                        generator.animator.avatar_colors[j], 
                        'excited', 
                        frame_count + i + j * 15
                    )
                    x = 150 + j * 250
                    y = height - 200 + int(30 * math.sin((frame_count + i + j * 10) * 0.2))
                    img.paste(avatar, (x, y), avatar)
                
                # Add celebration particles
                for k in range(50):
                    particle_x = random.randint(0, width)
                    particle_y = random.randint(0, height)
                    particle_color = random.choice(point_colors)
                    particle_size = random.randint(2, 8)
                    draw.ellipse([particle_x, particle_y, 
                                particle_x + particle_size, particle_y + particle_size],
                               fill=particle_color)
                
                img.save(os.path.join(frames_dir, f"frame_{frame_count:05d}.jpg"), quality=95)
                frame_count += 1
            
            # Use FFMPEG to create video with better quality
            ffmpeg_cmd = [
                'ffmpeg',
                '-y',  # Overwrite output file
                '-framerate', '24',
                '-i', os.path.join(frames_dir, 'frame_%05d.jpg'),
                '-c:v', 'libx264',
                '-preset', 'medium',
                '-crf', '18',  # Better quality
                '-pix_fmt', 'yuv420p',
                '-movflags', '+faststart',  # Enable streaming
                output_path
            ]
            
            # Run ffmpeg
            result = subprocess.run(ffmpeg_cmd, capture_output=True, text=True)
            if result.returncode != 0:
                return jsonify({
                    'error': f'FFMPEG error: {result.stderr}',
                    'stdout': result.stdout
                }), 500
            
            # Check if the video was created
            if not os.path.exists(output_path):
                return jsonify({'error': 'Failed to create video file'}), 500
            
            # Return the video file
            return send_file(output_path, as_attachment=True, 
                            download_name='key_points_video.mp4',
                            mimetype='video/mp4')
            
        finally:
            # Clean up temporary files (with delay to ensure file is sent)
            import shutil
            import threading
            
            def cleanup_later():
                import time
                time.sleep(5)  # Wait 5 seconds before cleanup
                try:
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    print(f"Error cleaning up: {e}")
            
            cleanup_thread = threading.Thread(target=cleanup_later)
            cleanup_thread.daemon = True
            cleanup_thread.start()
                
    except Exception as e:
        return jsonify({
            'error': f'Error generating enhanced video: {str(e)}',
            'traceback': traceback.format_exc()
        }), 500
                
## A new feature for mcq generation
def generate_mcqs_from_content(text, num_questions=5):
    """Generate MCQs using Gemini from the given content."""
    prompt = f"""
Based on the following research content, generate {num_questions} multiple-choice questions (MCQs) for students. 
Each MCQ should have:
- A question (1-2 lines)
- Four options (A, B, C, D)
- The correct answer (just the letter)
- A 1-line explanation for the correct answer

Format:
Q1: <question>
A) <option A>
B) <option B>
C) <option C>
D) <option D>
Answer: <A/B/C/D>
Explanation: <1-line explanation>

Content:
{text[:3000]}
"""
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        raise Exception(f"Error generating MCQs: {e}")

@app.route('/generate-mcq', methods=['POST'])
def handle_generate_mcq():
    try:
        data = request.json
        text = data.get('text')
        num_questions = int(data.get('num_questions', 5))

        if not text or len(text.strip()) < 30:
            return jsonify({'error': 'Please provide sufficient content for MCQ generation.'}), 400

        mcq_text = generate_mcqs_from_content(text, num_questions)
        # Optionally, parse MCQs into a structured format for frontend
        mcqs = []
        import re
        pattern = re.compile(
        r"Q\d+:(.*?)\nA\)(.*?)\nB\)(.*?)\nC\)(.*?)\nD\)(.*?)\nAnswer:\s*([A-D])\nExplanation:(.*?)(?:\n|$)",
        re.DOTALL
        )
        
        for match in pattern.finditer(mcq_text):
            mcqs.append({
                "question": match.group(1).strip(),
                "options": [
                    match.group(2).strip(),
                    match.group(3).strip(),
                    match.group(4).strip(),
                    match.group(5).strip()
                ],
                "answer": match.group(6).strip(),
                "explanation":match.group(7).strip()
            })

        return jsonify({
            "mcqs": mcqs,
            "raw": mcq_text
        }), 200
    except Exception as e:
        return jsonify({'error': f'Error generating MCQs: {str(e)}'}), 500

                                        
if __name__ == '__main__':
    os.makedirs('templates', exist_ok=True)
    port = int(os.environ.get("PORT", 5000))
    print(f"Starting server on port {port}")
    app.run(host="0.0.0.0", port=port, debug=False, threaded=True)        
